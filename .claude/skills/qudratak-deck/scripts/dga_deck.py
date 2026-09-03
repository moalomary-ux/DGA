# -*- coding: utf-8 -*-
"""dga_deck.py - reusable helpers for building DGA / Qudratak (قدراتك) decks with python-pptx.

Design system encoded here (colors, fonts, blocks) comes from the DGA brand identity and the
Qudratak progress-report template. Import with `from dga_deck import *` and build on top of a
template-derived .pptx (see build_qudratak_report.py for a full example).
"""
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.oxml.ns import qn

import os
SRC = os.environ.get('DECK_SRC', 'base.pptx')
OUT = os.environ.get('DECK_OUT', 'out.pptx')
MEDIA = os.environ.get('DECK_MEDIA', 'work/ppt/media/')
ASSETS = os.environ.get('DECK_ASSETS', '.')
DECK_TITLE = os.environ.get('DECK_TITLE', 'برنامج القدرات الرقمية «قدراتك»   |   تقرير سير الأعمال الربعي – يونيو – أغسطس 2026')

# ---- brand ---------------------------------------------------------------
NAVY = '2A206A'; PURPLE = '7C32C9'; PURPLE_DK = '5D2597'; TEAL = '00ABAF'; GREEN = '1CC182'
WHITE = 'FFFFFF'; TXT = '333333'; MUTED = '6E6E86'; CARD = 'F4F3F9'; LINE = 'E1E0EC'
AMBER = 'F0A322'; RED = 'E0525E'; GREEN_DK = '119C68'
F_BOLD = 'DiodrumArabic-Bold'; F_SEMI = 'DiodrumArabic-Semibold'
F_MED = 'DiodrumArabic-Medium'; F_REG = 'DiodrumArabic-Regular'

ICON_PROGRAM = MEDIA + 'image20.png'
ICON_WORKSHOP = MEDIA + 'image22.png'
LOGO_MCIT = os.path.join(ASSETS, 'mcit_logo.png')
LOGO_AWARD = MEDIA + 'image28.png'
LOGO_GCLOUD = MEDIA + 'image24.png'
LOGO_UNI_GREEN = MEDIA + 'image25.png'
LOGO_KSU = MEDIA + 'image26.png'

SW, SH = 13.333, 7.5
XL, XR = 0.45, 12.9          # content left / right edges
CW = XR - XL                 # content width


# ---- low-level helpers ---------------------------------------------------
def _rm(el):
    el.getparent().remove(el)


def strip_style(shape):
    st = shape._element.find(qn('p:style'))
    if st is not None:
        _rm(st)


def set_fill(shape, color, alpha=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    if alpha is not None:
        clr = shape.fill._xPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        a = etree.SubElement(clr, qn('a:alpha')); a.set('val', str(int(alpha * 1000)))


def set_grad(shape, c1, c2, ang=2700000, pos2=90000):
    spPr = shape._element.spPr
    for tag in ('a:solidFill', 'a:noFill', 'a:gradFill'):
        e = spPr.find(qn(tag))
        if e is not None:
            _rm(e)
    g = etree.Element(qn('a:gradFill'))
    gs = etree.SubElement(g, qn('a:gsLst'))
    for pos, c in ((0, c1), (pos2, c2)):
        s = etree.SubElement(gs, qn('a:gs')); s.set('pos', str(pos))
        cc = etree.SubElement(s, qn('a:srgbClr')); cc.set('val', c)
    lin = etree.SubElement(g, qn('a:lin')); lin.set('ang', str(ang)); lin.set('scaled', '1')
    # insert after prstGeom
    geom = spPr.find(qn('a:prstGeom'))
    geom.addnext(g)


def no_line(shape):
    shape.line.fill.background()


def set_line(shape, color, w=0.75):
    shape.line.color.rgb = RGBColor.from_string(color)
    shape.line.width = Pt(w)


def add_shadow(shape, blur=76200, dist=23812, alpha=12000):
    spPr = shape._element.spPr
    old = spPr.find(qn('a:effectLst'))
    if old is not None:
        _rm(old)
    eff = etree.SubElement(spPr, qn('a:effectLst'))
    sh = etree.SubElement(eff, qn('a:outerShdw'))
    sh.set('blurRad', str(blur)); sh.set('dist', str(dist)); sh.set('dir', '5400000'); sh.set('rotWithShape', '0')
    c = etree.SubElement(sh, qn('a:srgbClr')); c.set('val', '24135F')
    a = etree.SubElement(c, qn('a:alpha')); a.set('val', str(alpha))


def rect(slide, x, y, w, h, fill=WHITE, radius=0.06, shadow=True, line=None, grad=None, alpha=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    strip_style(shp)
    if radius is not None:
        shp.adjustments[0] = radius / min(w, h) if min(w, h) > 0 else 0.05
    if grad:
        set_grad(shp, *grad)
    else:
        set_fill(shp, fill, alpha)
    if line:
        set_line(shp, line)
    else:
        no_line(shp)
    if shadow:
        add_shadow(shp)
    shp.text_frame.text = ''
    return shp


def oval(slide, x, y, d, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    strip_style(shp)
    set_fill(shp, fill)
    if line:
        set_line(shp, line, 1.25)
    else:
        no_line(shp)
    return shp


def style_run(r, size=12, font=F_MED, color=TXT, bold=False, lang='ar-SA'):
    rPr = r._r.get_or_add_rPr()
    for ch in list(rPr):
        rPr.remove(ch)
    rPr.set('lang', lang); rPr.set('sz', str(int(round(size * 100))))
    rPr.set('b', '1' if bold else '0')
    rPr.set('dirty', '0')
    sf = etree.SubElement(rPr, qn('a:solidFill'))
    c = etree.SubElement(sf, qn('a:srgbClr')); c.set('val', color)
    for t in ('a:latin', 'a:cs'):
        e = etree.SubElement(rPr, qn(t)); e.set('typeface', font); e.set('pitchFamily', '2'); e.set('charset', '-78')


def style_par(p, align='r', rtl=True, line=None, after=0, before=0):
    pPr = p._p.get_or_add_pPr()
    for ch in list(pPr):
        pPr.remove(ch)
    pPr.set('algn', align)
    pPr.set('rtl', '1' if rtl else '0')
    if line:
        ls = etree.SubElement(pPr, qn('a:lnSpc')); sp = etree.SubElement(ls, qn('a:spcPct')); sp.set('val', str(int(line * 1000)))
    if before:
        sb = etree.SubElement(pPr, qn('a:spcBef')); sp = etree.SubElement(sb, qn('a:spcPts')); sp.set('val', str(int(before * 100)))
    if after:
        sa = etree.SubElement(pPr, qn('a:spcAft')); sp = etree.SubElement(sa, qn('a:spcPts')); sp.set('val', str(int(after * 100)))
    # remove endParaRPr fonts that could override
    return pPr


def add_bullet(p, color=PURPLE, char='•', indent=0.18):
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', str(int(Inches(indent)))); pPr.set('indent', str(-int(Inches(indent))))
    bc = etree.SubElement(pPr, qn('a:buClr')); c = etree.SubElement(bc, qn('a:srgbClr')); c.set('val', color)
    bf = etree.SubElement(pPr, qn('a:buFont')); bf.set('typeface', 'Arial')
    b = etree.SubElement(pPr, qn('a:buChar')); b.set('char', char)


def text(slide, x, y, w, h, paras, size=12, font=F_MED, color=TXT, bold=False, align='r',
         anchor='t', line=None, after=0, wrap=True, margin=0, autofit=False, bullets=None, bullet_color=PURPLE):
    """paras: str | list of (str | list of runs). run = str | (str, dict)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = {'t': MSO_ANCHOR.TOP, 'm': MSO_ANCHOR.MIDDLE, 'b': MSO_ANCHOR.BOTTOM}[anchor]
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if autofit:
        etree.SubElement(bodyPr, qn('a:spAutoFit'))
    else:
        etree.SubElement(bodyPr, qn('a:noAutofit'))
    if isinstance(paras, str):
        paras = [paras]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        style_par(p, align=align, line=line, after=after)
        if bullets:
            add_bullet(p, color=bullet_color)
        runs = [para] if isinstance(para, str) else para
        for run in runs:
            if isinstance(run, str):
                txt, opt = run, {}
            else:
                txt, opt = run
            r = p.add_run(); r.text = txt
            style_run(r, size=opt.get('size', size), font=opt.get('font', font), color=opt.get('color', color),
                      bold=opt.get('bold', bold))
    return tb


def picture(slide, path, x, y, w=None, h=None):
    kw = {}
    if w is not None:
        kw['width'] = Inches(w)
    if h is not None:
        kw['height'] = Inches(h)
    return slide.shapes.add_picture(path, Inches(x), Inches(y), **kw)


def line_h(slide, x1, x2, y, color=LINE, w=0.75, grad=None):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y))
    strip_style(c)
    if grad:
        spPr = c._element.spPr
        ln = spPr.find(qn('a:ln'))
        if ln is not None:
            _rm(ln)
        ln = etree.SubElement(spPr, qn('a:ln')); ln.set('w', str(int(Pt(w))))
        g = etree.SubElement(ln, qn('a:gradFill')); gs = etree.SubElement(g, qn('a:gsLst'))
        for pos, col in ((0, grad[0]), (100000, grad[1])):
            s = etree.SubElement(gs, qn('a:gs')); s.set('pos', str(pos))
            cc = etree.SubElement(s, qn('a:srgbClr')); cc.set('val', col)
        lin = etree.SubElement(g, qn('a:lin')); lin.set('ang', '0'); lin.set('scaled', '1')
    else:
        c.line.color.rgb = RGBColor.from_string(color); c.line.width = Pt(w)
    return c


def line_v(slide, x, y1, y2, color=LINE, w=0.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y1), Inches(x), Inches(y2))
    strip_style(c)
    c.line.color.rgb = RGBColor.from_string(color); c.line.width = Pt(w)
    return c


def chip(slide, x, y, w, h, label, fill=None, color=WHITE, size=10.5, font=F_SEMI, alpha=None, line=None):
    r = rect(slide, x, y, w, h, fill=fill or WHITE, radius=h / 2, shadow=False, alpha=alpha, line=line)
    tf = r.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08); tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; style_par(p, align='ctr')
    run = p.add_run(); run.text = label; style_run(run, size=size, font=font, color=color)
    return r


def status_pill(slide, x, y, w, label, kind='ok'):
    col = {'ok': GREEN, 'warn': AMBER, 'risk': RED, 'info': TEAL}[kind]
    r = rect(slide, x, y, w, 0.34, fill=col, radius=0.17, shadow=False, alpha=14)
    oval(slide, x + w - 0.27, y + 0.09, 0.16, col)
    tf = r.text_frame; tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.34); tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; style_par(p, align='r')
    run = p.add_run(); run.text = label; style_run(run, size=10.5, font=F_SEMI, color={'ok': GREEN_DK, 'warn': 'B36E00', 'risk': 'B0323C', 'info': '007A7D'}[kind])
    return r


def set_notes(slide, txt):
    slide.notes_slide.notes_text_frame.text = txt


def set_shape_text(shape, txt, keep_runs=1):
    """Replace text while keeping first-run formatting."""
    p = shape.text_frame.paragraphs[0]
    runs = p.runs
    runs[0].text = txt
    for r in runs[1:]:
        _rm(r._r)
    for extra in shape.text_frame.paragraphs[1:]:
        _rm(extra._p)


def clear_content(slide, keep_types=('sldNum',)):
    for shp in list(slide.shapes):
        keep = False
        if shp.is_placeholder:
            t = shp._element.find('.//' + qn('p:ph')).get('type')
            if t in keep_types:
                keep = True
        if not keep:
            _rm(shp._element)


def find(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    raise KeyError(name)


def ph_type(shape):
    ph = shape._element.find('.//' + qn('p:ph'))
    return ph.get('type') if ph is not None else None


def header_page_number(slide, y=0.17, h=0.6, size=12, color=NAVY):
    """Move the slide-number placeholder into the header row (left side, same baseline as the title)."""
    for shp in slide.shapes:
        if shp.is_placeholder and ph_type(shp) == 'sldNum':
            shp.left = Inches(XL); shp.top = Inches(y); shp.width = Inches(0.9); shp.height = Inches(h)
            txBody = shp._element.txBody
            bodyPr = txBody.find(qn('a:bodyPr')); bodyPr.set('anchor', 'ctr'); bodyPr.set('lIns', '0')
            for para in txBody.findall(qn('a:p')):
                pPr = para.find(qn('a:pPr'))
                if pPr is None:
                    pPr = etree.Element(qn('a:pPr')); para.insert(0, pPr)
                pPr.set('algn', 'l')
                for fld in para.findall(qn('a:fld')):
                    rPr = fld.find(qn('a:rPr'))
                    if rPr is None:
                        rPr = etree.Element(qn('a:rPr')); fld.insert(0, rPr)
                    for ch in list(rPr):
                        rPr.remove(ch)
                    rPr.set('lang', 'en-US'); rPr.set('sz', str(size * 100)); rPr.set('b', '1')
                    sf = etree.SubElement(rPr, qn('a:solidFill')); c = etree.SubElement(sf, qn('a:srgbClr')); c.set('val', color)
                    for t in ('a:latin', 'a:cs'):
                        e = etree.SubElement(rPr, qn(t)); e.set('typeface', F_SEMI)
            return shp


def remove_footers(prs):
    """Drop the classification footer placeholder ("Restricted - مقيد") from every slide."""
    for slide in prs.slides:
        for shp in list(slide.shapes):
            if shp.is_placeholder and ph_type(shp) == 'ftr':
                _rm(shp._element)


def footer(slide, title=None):
    """DNA footer on content slides: page number bottom-left (under the layout's small green dash),
    deck title bottom-right next to the DGA watermark. The classification label is intentionally absent."""
    for shp in slide.shapes:
        if shp.is_placeholder and ph_type(shp) == 'sldNum':
            shp.left = Inches(0.42); shp.top = Inches(7.06); shp.width = Inches(0.8); shp.height = Inches(0.26)
            txBody = shp._element.txBody
            bodyPr = txBody.find(qn('a:bodyPr')); bodyPr.set('anchor', 'ctr'); bodyPr.set('lIns', '0'); bodyPr.set('tIns', '0'); bodyPr.set('bIns', '0')
            for para in txBody.findall(qn('a:p')):
                pPr = para.find(qn('a:pPr'))
                if pPr is None:
                    pPr = etree.Element(qn('a:pPr')); para.insert(0, pPr)
                pPr.set('algn', 'l')
                for fld in para.findall(qn('a:fld')):
                    rPr = fld.find(qn('a:rPr'))
                    if rPr is None:
                        rPr = etree.Element(qn('a:rPr')); fld.insert(0, rPr)
                    for ch in list(rPr):
                        rPr.remove(ch)
                    rPr.set('lang', 'en-US'); rPr.set('sz', '1050'); rPr.set('b', '1')
                    sf = etree.SubElement(rPr, qn('a:solidFill')); c = etree.SubElement(sf, qn('a:srgbClr')); c.set('val', NAVY)
                    for t in ('a:latin', 'a:cs'):
                        e = etree.SubElement(rPr, qn(t)); e.set('typeface', F_SEMI)
    text(slide, 7.2, 6.84, 12.8 - 7.2, 0.26, title or DECK_TITLE, size=9, font=F_REG, color=MUTED, anchor='m')


# ---- composite blocks ----------------------------------------------------
def content_frame(slide, num, title, lead=None):
    """Title (number in purple + text in navy), gradient rule, optional teal lead line."""
    clear_content(slide)
    runs = []
    if num:
        runs.append((num + '  ', {'color': PURPLE_DK}))
    runs.append((title, {}))
    text(slide, XL, 0.17, CW, 0.6, [runs], size=24, font=F_BOLD, color=NAVY, anchor='m')
    footer(slide)
    line_h(slide, XL, XR + 0.04, 0.87, grad=(TEAL, PURPLE), w=2.25)
    if lead:
        text(slide, XL, 1.0, CW, 0.36, lead, size=13, font=F_SEMI, color=TEAL, anchor='m')


def hero_band(slide, x, y, w, h):
    return rect(slide, x, y, w, h, radius=0.12, grad=('1D164B', '3B2D8A'), shadow=True)


def kpi_on_dark(slide, x, y, w, number, label, sub=None, num_size=34):
    text(slide, x, y, w, 0.62, number, size=num_size, font=F_BOLD, color=WHITE, anchor='b')
    text(slide, x, y + 0.64, w, 0.3, label, size=12.5, font=F_SEMI, color=GREEN, anchor='t')
    if sub:
        text(slide, x, y + 0.94, w, 0.28, sub, size=10.5, font=F_REG, color='D9D6F0', anchor='t')


def kpi_card(slide, x, y, w, h, number, label, sub=None, num_color=PURPLE, icon=None):
    rect(slide, x, y, w, h, fill=WHITE, radius=0.1, shadow=True, line=LINE)
    text(slide, x + 0.2, y + 0.18, w - 0.4, 0.62, number, size=30, font=F_BOLD, color=num_color, anchor='m')
    text(slide, x + 0.2, y + 0.82, w - 0.4, 0.3, label, size=12, font=F_SEMI, color=NAVY)
    if sub:
        text(slide, x + 0.2, y + 1.1, w - 0.4, 0.3, sub, size=10, font=F_REG, color=MUTED)
    if icon:
        picture(slide, icon, x + 0.18, y + 0.2, h=0.5)


def card(slide, x, y, w, h, title=None, title_color=PURPLE, fill=WHITE, line=LINE):
    rect(slide, x, y, w, h, fill=fill, radius=0.1, shadow=True, line=line)
    if title:
        text(slide, x + 0.22, y + 0.16, w - 0.44, 0.38, title, size=13.5, font=F_BOLD, color=title_color, anchor='m')
        line_h(slide, x + 0.22, x + w - 0.22, y + 0.6, color=LINE, w=0.75)


def table(slide, x, y, col_w, rows, header, row_h=0.5, head_h=0.42, size=10.5, head_fill=TEAL, align_cols=None,
          zebra=True, first_bold=True, pad_right=None):
    """col_w and header listed in RTL logical order (first = rightmost)."""
    ncol = len(col_w); nrow = len(rows) + 1
    w = sum(col_w)
    h = head_h + row_h * len(rows)
    gf = slide.shapes.add_table(nrow, ncol, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    sid = tblPr.find(qn('a:tableStyleId'))
    if sid is not None:
        _rm(sid)
    # visual order: reverse (rightmost logical first)
    vis_w = list(reversed(col_w)); vis_head = list(reversed(header))
    for i, cw in enumerate(vis_w):
        tbl.columns[i].width = Inches(cw)
    tbl.rows[0].height = Inches(head_h)
    for r in range(1, nrow):
        tbl.rows[r].height = Inches(row_h)

    def fill_cell(cell, txt, fnt, col, sz, bold, fillc, al='r', pr=0.08):
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor.from_string(fillc)
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(pr); cell.margin_top = cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        paras = txt if isinstance(txt, list) else [txt]
        for i, ptxt in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            style_par(p, align=al, line=110)
            run = p.add_run(); run.text = ptxt
            style_run(run, size=sz, font=fnt, color=col, bold=bold)
        # borders
        tcPr = cell._tc.get_or_add_tcPr()
        for edge in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
            ln = etree.SubElement(tcPr, qn(edge)); ln.set('w', '6350')
            sf = etree.SubElement(ln, qn('a:solidFill')); c = etree.SubElement(sf, qn('a:srgbClr')); c.set('val', 'FFFFFF')
        # tcPr children order: lnL lnR lnT lnB ... then fill; move solidFill after borders
        fillel = tcPr.find(qn('a:solidFill'))
        if fillel is not None:
            tcPr.remove(fillel); tcPr.append(fillel)

    for c in range(ncol):
        fill_cell(tbl.cell(0, c), vis_head[c], F_SEMI, WHITE, size, False, head_fill, 'ctr')
    for r, row in enumerate(rows, start=1):
        vis_row = list(reversed(row))
        fillc = (CARD if (r % 2 == 1) else WHITE) if zebra else WHITE
        for c in range(ncol):
            logical = ncol - 1 - c
            al = 'r'
            if align_cols and align_cols[logical]:
                al = align_cols[logical]
            bold = first_bold and logical == 0
            pr = (pad_right or {}).get(logical, 0.08)
            fill_cell(tbl.cell(r, c), vis_row[c], F_SEMI if bold else F_MED, NAVY if bold else TXT, size, False, fillc, al, pr)
    return gf




__all__ = [n for n in dir() if not n.startswith('__')]
