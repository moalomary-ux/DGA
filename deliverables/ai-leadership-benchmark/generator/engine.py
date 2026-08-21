# -*- coding: utf-8 -*-
"""Deck engine: renders a bilingual content model into an A4-portrait PPTX
following the Rayan-approved DGA identity. All shapes/tables stay native
and editable; only the cover background and one contextual banner are images."""
import copy
from PIL import ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import theme as T

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# ---------------------------------------------------------------- measurement
_font_cache = {}
def _pil(stylekey, size):
    key = (stylekey, size)
    if key not in _font_cache:
        _font_cache[key] = ImageFont.truetype(T.FONTS[stylekey][2], int(size * 8))
    return _font_cache[key]

def text_w(text, stylekey, size, spc=0):
    """width of text in inches (spc = letterspacing in pt per char)"""
    if not text:
        return 0.0
    f = _pil(stylekey, size)
    w = f.getlength(text) / 8.0 / 72.0
    if spc:
        w += len(text) * spc / 72.0
    return w

def wrap(text, stylekey, size, width_in, spc=0):
    lines = []
    for hard in str(text).split('\n'):
        words = hard.split()
        if not words:
            lines.append('')
            continue
        cur = words[0]
        for w in words[1:]:
            if text_w(cur + ' ' + w, stylekey, size, spc) <= width_in:
                cur += ' ' + w
            else:
                lines.append(cur)
                cur = w
        lines.append(cur)
    return lines

def n_lines(text, stylekey, size, width_in, spc=0):
    return max(1, len(wrap(text, stylekey, size, width_in, spc)))

def para_h(text, stylekey, size, width_in, spacing=1.16, spc=0):
    """spacing is an ABSOLUTE leading ratio: line height = size * spacing pt."""
    return n_lines(text, stylekey, size, width_in, spc) * size * spacing / 72.0

# ---------------------------------------------------------------- low level
def _rgb(c):
    return RGBColor(*c)

def set_run_font(run, stylekey, size, color):
    fam, bold, _ = T.FONTS[stylekey]
    run.font.name = fam
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    rPr = run._r.get_or_add_rPr()
    for tag in ('ea', 'cs'):
        el = rPr.find('{%s}%s' % (A_NS, tag))
        if el is None:
            el = etree.SubElement(rPr, '{%s}%s' % (A_NS, tag))
        el.set('typeface', fam)

def set_rtl_para(p):
    p._p.get_or_add_pPr().set('rtl', '1')

def set_spc(run, spc_pt):
    run._r.get_or_add_rPr().set('spc', str(int(spc_pt * 100)))

def no_line(shape):
    shape.line.fill.background()

def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)

def gradient(shape, c1, c2, angle=0):
    shape.fill.gradient()
    stops = shape.fill.gradient_stops
    stops[0].color.rgb = _rgb(c1)
    stops[1].color.rgb = _rgb(c2)
    try:
        shape.fill.gradient_angle = angle
    except Exception:
        pass

def soft_shadow(shape):
    spPr = shape._element.spPr
    el = spPr.find(qn('a:effectLst'))
    if el is None:
        el = etree.SubElement(spPr, qn('a:effectLst'))
    sh = etree.SubElement(el, qn('a:outerShdw'))
    sh.set('blurRad', '90000'); sh.set('dist', '30000')
    sh.set('dir', '5400000'); sh.set('rotWithShape', '0')
    clr = etree.SubElement(sh, qn('a:srgbClr')); clr.set('val', '2A206A')
    a = etree.SubElement(clr, qn('a:alpha')); a.set('val', '10000')

def rounded(slide, x, y, w, h, radius=0.07):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        sp.adjustments[0] = max(0.02, min(0.5, radius / min(w, h)))
    except Exception:
        pass
    return sp

def rect(slide, x, y, w, h):
    return slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))

# ---------------------------------------------------------------- text block
def add_text(slide, x, y, w, paras, rtl=False, anchor=None, wrap_on=True):
    """paras: list of dicts:
       {text | runs:[{text, style, size, color, spc?, link?}], style, size, color,
        align('r'/'l'/'c'/'j'/None=start), spacing, after, before}
       returns height in inches (measured)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = wrap_on
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor == 'middle':
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    total_h = 0.0
    for i, spec in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        align = spec.get('align')
        if align is None:
            align = 'r' if rtl else 'l'
        p.alignment = {'r': PP_ALIGN.RIGHT, 'l': PP_ALIGN.LEFT, 'c': PP_ALIGN.CENTER, 'j': PP_ALIGN.JUSTIFY}[align]
        if rtl and not spec.get('ltr'):
            set_rtl_para(p)
        spacing = spec.get('spacing', 1.16)
        if spec.get('after'):
            p.space_after = Pt(spec['after'])
        if spec.get('before'):
            p.space_before = Pt(spec['before'])
        runs = spec.get('runs') or [{'text': spec.get('text', ''), 'style': spec.get('style', 'reg'),
                                     'size': spec.get('size', 10), 'color': spec.get('color', T.BODY),
                                     'spc': spec.get('spc'), 'link': spec.get('link')}]
        joined = ''.join(r['text'] for r in runs)
        style0 = runs[0].get('style', 'reg'); size0 = max(r.get('size', 10) for r in runs)
        p.line_spacing = Pt(size0 * spacing)
        for rs in runs:
            r = p.add_run()
            r.text = rs['text']
            set_run_font(r, rs.get('style', 'reg'), rs.get('size', 10), rs.get('color', T.BODY))
            if rs.get('spc'):
                set_spc(r, rs['spc'])
            if rs.get('link'):
                r.hyperlink.address = rs['link']
        lh = size0 * spacing / 72.0
        spc0 = runs[0].get('spc') or 0
        nl = n_lines(joined, style0, size0, w, spc0) if wrap_on else 1
        total_h += nl * lh
        total_h += spec.get('after', 0) / 72.0 + spec.get('before', 0) / 72.0
    tb.height = Inches(max(0.18, total_h + 0.04))
    return total_h

# ---------------------------------------------------------------- footer/header
def footer(slide, page_no, doc_title, rtl):
    ln = rect(slide, T.MARGIN, T.FOOTER_RULE_Y, T.CONTENT_W, 0.008)
    solid(ln, T.CARD_BORDER); no_line(ln)
    if rtl:
        add_text(slide, T.PAGE_W - T.MARGIN - 4.0, T.FOOTER_TEXT_Y, 4.0,
                 [{'text': doc_title, 'style': 'reg', 'size': 6.5, 'color': T.FOOT_GRAY, 'align': 'r'}], rtl=True)
        add_text(slide, T.MARGIN, T.FOOTER_TEXT_Y, 0.6,
                 [{'text': str(page_no), 'style': 'med', 'size': 7, 'color': T.FOOT_GRAY, 'align': 'l'}])
    else:
        add_text(slide, T.MARGIN, T.FOOTER_TEXT_Y, 4.0,
                 [{'text': doc_title, 'style': 'reg', 'size': 6.5, 'color': T.FOOT_GRAY, 'align': 'l'}])
        add_text(slide, T.PAGE_W - T.MARGIN - 0.6, T.FOOTER_TEXT_Y, 0.6,
                 [{'text': str(page_no), 'style': 'med', 'size': 7, 'color': T.FOOT_GRAY, 'align': 'r'}])

def grad_tick(slide, x, y, rtl, w=0.55):
    tx = (T.PAGE_W - T.MARGIN - w) if rtl else x
    sp = rect(slide, tx, y, w, 0.045)
    gradient(sp, T.PURPLE, T.TEAL, 0)
    no_line(sp)

# ---------------------------------------------------------------- blocks
class Ctx:
    def __init__(self, slide, rtl, warns, page_no):
        self.slide = slide
        self.rtl = rtl
        self.warns = warns
        self.page = page_no

def _tick_color(c):
    return {'purple': T.PURPLE, 'teal': T.TEAL, 'green': T.GREEN, 'blue': T.BLUE,
            'navy': T.NAVY, 'red': T.RED}.get(c, c if isinstance(c, tuple) else T.PURPLE)

def blk_kicker(ctx, b, x, y, w):
    spc = 0 if ctx.rtl else 0.9
    h = add_text(ctx.slide, x, y, w, [{'text': b['text'], 'style': 'semi', 'size': 10.5,
                                       'color': T.PURPLE, 'spc': spc}], rtl=ctx.rtl)
    return y + h + 0.08

def blk_h1(ctx, b, x, y, w):
    size = b.get('size', 20)
    h = add_text(ctx.slide, x, y, w, [{'text': b['text'], 'style': 'bold', 'size': size,
                                       'color': T.INK, 'spacing': 1.14}], rtl=ctx.rtl)
    return y + h + 0.10

def blk_lede(ctx, b, x, y, w):
    h = add_text(ctx.slide, x, y, w, [{'text': b['text'], 'style': 'reg', 'size': 9.5,
                                       'color': T.GRAY, 'spacing': 1.40}], rtl=ctx.rtl)
    grad_tick(ctx.slide, x, y + h + 0.10, ctx.rtl)
    return y + h + 0.30

def blk_body(ctx, b, x, y, w):
    size = b.get('size', 9.5)
    style = b.get('style', 'reg')
    align = b.get('align', 'j' if len(str(b['text'])) > 200 else None)
    h = add_text(ctx.slide, x, y, w, [{'text': b['text'], 'style': style, 'size': size,
                                       'color': b.get('color', T.BODY), 'spacing': 1.42,
                                       'align': align}], rtl=ctx.rtl)
    return y + h + b.get('after', 0.14)

def blk_note(ctx, b, x, y, w):
    h = add_text(ctx.slide, x, y, w, [{'text': b['text'], 'style': 'reg', 'size': 8,
                                       'color': T.GRAY, 'spacing': 1.40}], rtl=ctx.rtl)
    return y + h + 0.12

def blk_gap(ctx, b, x, y, w):
    return y + b.get('h', 0.15)

def blk_subhead(ctx, b, x, y, w):
    color = _tick_color(b.get('color', 'purple'))
    spc = 0 if ctx.rtl else 0.7
    h = add_text(ctx.slide, x, y, w, [{'text': b['text'], 'style': 'semi', 'size': b.get('size', 11.5),
                                       'color': color, 'spc': spc}], rtl=ctx.rtl)
    return y + h + 0.09

def _card(ctx, x, y, w, h, item, title_size, body_size):
    card = rounded(ctx.slide, x, y, w, h, 0.09)
    solid(card, T.WHITE)
    card.line.color.rgb = _rgb(T.CARD_BORDER)
    card.line.width = Pt(0.75)
    soft_shadow(card)
    pad = 0.16
    tick_c = _tick_color(item.get('tick', 'purple'))
    tx = (x + w - pad - 0.38) if ctx.rtl else (x + pad)
    tick = rect(ctx.slide, tx, y + pad, 0.38, 0.05)
    solid(tick, tick_c); no_line(tick)
    ty = y + pad + 0.13
    tcol = item.get('title_color', tick_c)
    th = add_text(ctx.slide, x + pad, ty, w - 2 * pad,
                  [{'text': item['title'], 'style': 'semi', 'size': title_size, 'color': tcol,
                    'spacing': 1.22}], rtl=ctx.rtl)
    by = ty + th + 0.07
    if item.get('body'):
        add_text(ctx.slide, x + pad, by, w - 2 * pad,
                 [{'text': item['body'], 'style': 'reg', 'size': body_size, 'color': T.BODY,
                   'spacing': 1.38, 'align': 'j' if (len(item['body']) > 120 and w - 2 * pad > 2.8) else None}], rtl=ctx.rtl)

def _card_h(ctx, w, item, title_size, body_size):
    pad = 0.16
    h = pad + 0.13
    h += para_h(item['title'], 'semi', title_size, w - 2 * pad, 1.22) + 0.07
    if item.get('body'):
        h += para_h(item['body'], 'reg', body_size, w - 2 * pad, 1.38)
    return h + pad

def blk_cards(ctx, b, x, y, w):
    cols = b.get('cols', len(b['items']))
    gap = 0.16
    cw = (w - (cols - 1) * gap) / cols
    title_size = b.get('title_size', 10.5)
    body_size = b.get('body_size', 8.5)
    items = b['items']
    rows = [items[i:i + cols] for i in range(0, len(items), cols)]
    for row in rows:
        rh = max(_card_h(ctx, cw, it, title_size, body_size) for it in row)
        rh = max(rh, b.get('min_h', 0))
        for i, it in enumerate(row):
            ci = (len(row) - 1 - i) if ctx.rtl else i
            cx = x + ci * (cw + gap)
            _card(ctx, cx, y, cw, rh, it, title_size, body_size)
        y += rh + gap
    return y + 0.04

def blk_stats(ctx, b, x, y, w):
    items = b['items']
    cols = b.get('cols', len(items))
    gap = 0.16
    cw = (w - (cols - 1) * gap) / cols
    pad = 0.16
    num_size = b.get('num_size', 24)
    heights = []
    for it in items:
        h = pad + num_size * 1.1 / 72 + 0.08
        h += para_h(it['label'], 'reg', 8.5, cw - 2 * pad, 1.35)
        if it.get('src'):
            h += 0.06 + para_h(it['src'], 'reg', 7, cw - 2 * pad, 1.25)
        heights.append(h + pad)
    rh = max(heights)
    for i, it in enumerate(items):
        ci = (len(items) - 1 - i) if ctx.rtl else i
        cx = x + ci * (cw + gap)
        card = rounded(ctx.slide, cx, y, cw, rh, 0.09)
        solid(card, T.WHITE); card.line.color.rgb = _rgb(T.CARD_BORDER); card.line.width = Pt(0.75)
        soft_shadow(card)
        ny = y + pad
        add_text(ctx.slide, cx + pad, ny, cw - 2 * pad,
                 [{'text': it['num'], 'style': 'bold', 'size': num_size,
                   'color': _tick_color(it.get('color', 'purple')), 'spacing': 1.05}], rtl=ctx.rtl)
        ny += num_size * 1.1 / 72 + 0.08
        h2 = add_text(ctx.slide, cx + pad, ny, cw - 2 * pad,
                      [{'text': it['label'], 'style': 'reg', 'size': 8.5, 'color': T.BODY,
                        'spacing': 1.35}], rtl=ctx.rtl)
        if it.get('src'):
            add_text(ctx.slide, cx + pad, ny + h2 + 0.06, cw - 2 * pad,
                     [{'text': it['src'], 'style': 'reg', 'size': 7, 'color': T.GRAY, 'spacing': 1.25}], rtl=ctx.rtl)
    return y + rh + 0.20

def blk_panel(ctx, b, x, y, w):
    pad = 0.18
    ih = pad
    tspc = 0 if ctx.rtl else 0.8
    if b.get('title'):
        ih += para_h(b['title'], 'semi', 10.5, w - 2 * pad, 1.15, tspc) + 0.06
    ih += para_h(b['body'], 'reg', b.get('size', 9), w - 2 * pad, 1.42) + pad + 0.05
    panel = rounded(ctx.slide, x, y, w, ih, 0.09)
    solid(panel, T.LAV)
    panel.line.color.rgb = _rgb(T.LAV_BORDER); panel.line.width = Pt(0.75)
    ty = y + pad
    if b.get('title'):
        spc = 0 if ctx.rtl else 0.8
        th = add_text(ctx.slide, x + pad, ty, w - 2 * pad,
                      [{'text': b['title'], 'style': 'semi', 'size': 10.5,
                        'color': _tick_color(b.get('title_color', 'purple')), 'spc': spc}], rtl=ctx.rtl)
        ty += th + 0.06
    add_text(ctx.slide, x + pad, ty, w - 2 * pad,
             [{'text': b['body'], 'style': 'reg', 'size': b.get('size', 9), 'color': T.BODY,
               'spacing': 1.42, 'align': 'j' if len(b['body']) > 150 else None}], rtl=ctx.rtl)
    return y + ih + 0.18

def blk_gpanel(ctx, b, x, y, w):
    pad = 0.20
    paras = b.get('paras') or [(b.get('title'), b['body'])]
    ih = pad
    tspc = 0 if ctx.rtl else 0.9
    for t, bd in paras:
        if t:
            ih += para_h(t, 'semi', 10.5, w - 2 * pad, 1.15, tspc) + 0.05
        ih += para_h(bd, 'reg', 9, w - 2 * pad, 1.42) + 0.12
    ih += pad - 0.06
    panel = rounded(ctx.slide, x, y, w, ih, 0.10)
    gradient(panel, T.PURPLE, T.NAVY, 20)
    no_line(panel)
    soft_shadow(panel)
    ty = y + pad
    for t, bd in paras:
        if t:
            spc = 0 if ctx.rtl else 0.9
            th = add_text(ctx.slide, x + pad, ty, w - 2 * pad,
                          [{'text': t, 'style': 'semi', 'size': 10.5, 'color': (0xD9, 0xF3, 0xEE),
                            'spc': spc}], rtl=ctx.rtl)
            ty += th + 0.05
        bh = add_text(ctx.slide, x + pad, ty, w - 2 * pad,
                      [{'text': bd, 'style': 'reg', 'size': 9, 'color': T.WHITE, 'spacing': 1.42,
                        'align': 'j' if len(bd) > 150 else None}], rtl=ctx.rtl)
        ty += bh + 0.12
    return y + ih + 0.18

def blk_transition(ctx, b, x, y, w):
    bar_w = 0.05
    th = para_h(b['text'], 'semi', 10, w - 0.25, 1.38)
    bx = (x + w - bar_w) if ctx.rtl else x
    bar = rect(ctx.slide, bx, y + 0.02, bar_w, th + 0.06)
    solid(bar, T.TEAL); no_line(bar)
    tx = x if ctx.rtl else x + 0.22
    add_text(ctx.slide, tx, y + 0.03, w - 0.25,
             [{'text': b['text'], 'style': 'semi', 'size': 10, 'color': T.INK, 'spacing': 1.38}],
             rtl=ctx.rtl)
    return y + th + 0.24

def blk_findings(ctx, b, x, y, w):
    num_w = 0.5
    for i, it in enumerate(b['items']):
        tw = w - num_w - 0.1
        tx = x if ctx.rtl else x + num_w + 0.1
        nx = (x + w - num_w) if ctx.rtl else x
        add_text(ctx.slide, nx, y - 0.02, num_w,
                 [{'text': str(it.get('n', i + 1)), 'style': 'bold', 'size': 17,
                   'color': _tick_color(it.get('color', 'purple')),
                   'align': 'l' if not ctx.rtl else 'r'}], rtl=ctx.rtl)
        th = add_text(ctx.slide, tx, y, tw,
                      [{'text': it['title'], 'style': 'semi', 'size': b.get('title_size', 10.5),
                        'color': _tick_color(it.get('color', 'purple')), 'spacing': 1.24}], rtl=ctx.rtl)
        bh = add_text(ctx.slide, tx, y + th + 0.05, tw,
                      [{'text': it['body'], 'style': 'reg', 'size': b.get('size', 9), 'color': T.BODY,
                        'spacing': 1.40, 'align': 'j' if len(it['body']) > 150 else None}], rtl=ctx.rtl)
        y += th + bh + 0.12
        if i < len(b['items']) - 1:
            ln = rect(ctx.slide, x, y, w, 0.008)
            solid(ln, T.CARD_BORDER); no_line(ln)
            y += 0.14
    return y + 0.08

def blk_bullets(ctx, b, x, y, w):
    size = b.get('size', 9)
    color = b.get('color', T.BODY)
    indent = 0.19
    for it in b['items']:
        bx = (x + w - indent) if ctx.rtl else x
        add_text(ctx.slide, bx if not ctx.rtl else x + w - 0.12, y + 0.012, 0.12,
                 [{'text': '•', 'style': 'semi', 'size': size, 'color': _tick_color(b.get('bullet_color', 'teal')),
                   'align': 'c'}]) if False else None
        # bullet dot
        dot_x = (x + w - 0.10) if ctx.rtl else x
        add_text(ctx.slide, dot_x, y - 0.005, 0.10,
                 [{'text': '•', 'style': 'semi', 'size': size + 1,
                   'color': _tick_color(b.get('bullet_color', 'teal')), 'align': 'c'}])
        tx = x if ctx.rtl else x + indent
        h = add_text(ctx.slide, tx, y, w - indent,
                     [{'text': it, 'style': 'reg', 'size': size, 'color': color, 'spacing': 1.38}],
                     rtl=ctx.rtl)
        y += h + b.get('item_gap', 0.055)
    return y + 0.08

def blk_chips(ctx, b, x, y, w):
    cols = b.get('cols', 4)
    gap = 0.14
    cw = (w - (cols - 1) * gap) / cols
    items = b['items']
    rows = [items[i:i + cols] for i in range(0, len(items), cols)]
    pad = 0.12
    for row in rows:
        heights = []
        for it in row:
            h = pad + 0.34 + para_h(it['label'], 'semi', 9, cw - 2 * pad, 1.28) + 0.04
            h += para_h(it['desc'], 'reg', 7.5, cw - 2 * pad, 1.32) + pad
            heights.append(h)
        rh = max(heights)
        for i, it in enumerate(row):
            ci = (len(row) - 1 - i) if ctx.rtl else i
            cx = x + ci * (cw + gap)
            card = rounded(ctx.slide, cx, y, cw, rh, 0.09)
            solid(card, T.WHITE); card.line.color.rgb = _rgb(T.CARD_BORDER); card.line.width = Pt(0.75)
            soft_shadow(card)
            cc = _tick_color(it.get('color', 'green'))
            badge = rounded(ctx.slide, (cx + cw - pad - 0.30) if ctx.rtl else cx + pad, y + pad, 0.30, 0.24, 0.06)
            solid(badge, cc); no_line(badge)
            tfb = badge.text_frame; tfb.word_wrap = False
            tfb.margin_left = tfb.margin_right = tfb.margin_top = tfb.margin_bottom = 0
            pb = tfb.paragraphs[0]; pb.alignment = PP_ALIGN.CENTER
            rb = pb.add_run(); rb.text = it['code']
            set_run_font(rb, 'semi', 10, T.WHITE)
            ly = y + pad + 0.34
            lh = add_text(ctx.slide, cx + pad, ly, cw - 2 * pad,
                          [{'text': it['label'], 'style': 'semi', 'size': 9, 'color': T.INK,
                            'spacing': 1.28}], rtl=ctx.rtl)
            add_text(ctx.slide, cx + pad, ly + lh + 0.04, cw - 2 * pad,
                     [{'text': it['desc'], 'style': 'reg', 'size': 7.5, 'color': T.GRAY,
                       'spacing': 1.32}], rtl=ctx.rtl)
        y += rh + gap
    return y + 0.05

def blk_pill(ctx, b, x, y, w):
    txt = b['text']
    size = b.get('size', 8.5)
    tw = text_w(txt, 'semi', size) + 0.30
    px = (x + w - tw) if ctx.rtl else x
    pill = rounded(ctx.slide, px, y, tw, 0.26, 0.13)
    fillc = b.get('fill', T.LAV)
    solid(pill, fillc)
    pill.line.color.rgb = _rgb(b.get('border', T.LAV_BORDER)); pill.line.width = Pt(0.75)
    tf = pill.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    if ctx.rtl:
        set_rtl_para(p)
    r = p.add_run(); r.text = txt
    set_run_font(r, 'semi', size, b.get('color', T.PURPLE))
    return y + 0.26 + 0.12

# ---------------------------------------------------------------- tables
def _set_cell(cell, text, stylekey, size, color, rtl, align=None, fill=None, anchor_mid=True):
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(fill)
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(T.WHITE)
    cell.margin_left = cell.margin_right = Inches(0.07)
    cell.margin_top = cell.margin_bottom = Inches(0.045)
    if anchor_mid:
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    parts = str(text).split('\n')
    for i, part in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = Pt(size * 1.28)
        p.alignment = {'r': PP_ALIGN.RIGHT, 'l': PP_ALIGN.LEFT, 'c': PP_ALIGN.CENTER}[
            align or ('r' if rtl else 'l')]
        if rtl:
            set_rtl_para(p)
        r = p.add_run(); r.text = part
        set_run_font(r, stylekey, size, color)

def blk_table(ctx, b, x, y, w):
    headers = b.get('headers')
    rows = b['rows']
    nrows = len(rows) + (1 if headers else 0)
    ncols = len(rows[0]) if rows else len(headers)
    widths = b.get('widths') or [1.0 / ncols] * ncols
    size = b.get('size', 8.5)
    hsize = b.get('hsize', size)
    gt_shape = ctx.slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(0.4))
    gt = gt_shape.table
    tblPr = gt._tbl.find('{%s}tblPr' % A_NS)
    if ctx.rtl:
        tblPr.set('rtl', '1')
    tblPr.set('firstRow', '1' if headers else '0')
    tblPr.set('bandRow', '0')
    # remove default style to control fills fully
    for el in gt._tbl.findall('{%s}tblPr/{%s}tableStyleId' % (A_NS, A_NS)):
        el.getparent().remove(el)
    for j, frac in enumerate(widths):
        gt.columns[j].width = Inches(w * frac)
    r0 = 0
    header_fill = b.get('header_fill', T.NAVY_DEEP)
    if headers:
        for j, htxt in enumerate(headers):
            _set_cell(gt.cell(0, j), htxt, 'semi', hsize, T.WHITE, ctx.rtl, fill=header_fill)
        r0 = 1
    col_styles = b.get('col_styles', {})  # {col: {'style','color','align'}}
    row_fills = b.get('row_fills', {})    # {row_idx: color}
    cell_fills = b.get('cell_fills', {})  # {(row,col): color}
    for i, row in enumerate(rows):
        default_fill = row_fills.get(i, T.ALT_ROW if (b.get('banded', True) and i % 2 == 1) else T.WHITE)
        for j, val in enumerate(row):
            cs = col_styles.get(j, {})
            fill = cell_fills.get((i, j), default_fill)
            _set_cell(gt.cell(r0 + i, j), val,
                      cs.get('style', 'reg'), cs.get('size', size),
                      cs.get('color', T.BODY), ctx.rtl, align=cs.get('align'), fill=fill)
    # measured row heights
    total_h = 0.0
    pad_h = 0.045 * 2 + 0.045
    if headers:
        hh = max(para_h(str(h), 'semi', hsize, w * widths[j] - 0.16, 1.28)
                 for j, h in enumerate(headers)) + pad_h
        gt.rows[0].height = Inches(hh)
        total_h += hh
    rtl_factor = 0.88 if ctx.rtl else 1.0
    for i, row in enumerate(rows):
        rh = 0
        for j, val in enumerate(row):
            cs = col_styles.get(j, {})
            cw = (w * widths[j] - 0.16) * rtl_factor
            hcell = 0
            for part in str(val).split('\n'):
                hcell += para_h(part, cs.get('style', 'reg'), cs.get('size', size), cw, 1.28)
            rh = max(rh, hcell)
        rh += pad_h + (0.02 if ctx.rtl else 0.0)
        rh = max(rh, b.get('min_row_h', 0.24))
        gt.rows[r0 + i].height = Inches(rh)
        total_h += rh
    return y + total_h + 0.16

# ---------------------------------------------------------------- special
def blk_image(ctx, b, x, y, w):
    ctx.slide.shapes.add_picture(b['path'], Inches(x), Inches(y), Inches(w), Inches(b['h']))
    return y + b['h'] + 0.15

def blk_columns(ctx, b, x, y, w):
    widths = b.get('widths', [0.5, 0.5])
    gap = b.get('gap', 0.22)
    total = sum(widths)
    avail = w - gap * (len(widths) - 1)
    xs = []
    cur = x
    for frac in widths:
        xs.append((cur, avail * frac / total))
        cur += avail * frac / total + gap
    if ctx.rtl:
        cols_geo = []
        cur = x
        for frac in reversed(widths):
            cols_geo.append((cur, avail * frac / total))
            cur += avail * frac / total + gap
        xs = list(reversed(cols_geo))
    max_y = y
    for (cx, cw), blocks in zip(xs, b['columns']):
        cy = render_blocks(ctx, blocks, cx, y, cw)
        max_y = max(max_y, cy)
    return max_y + b.get('after', 0.05)

def blk_levelband(ctx, b, x, y, w):
    """Colored level header band with decision badge."""
    lc = T.LEVELS[b['level']]
    band = rounded(ctx.slide, x, y, w, 0.62, 0.09)
    gradient(band, lc, T.NAVY if b['level'] != 2 else T.PURPLE, 0)
    no_line(band)
    pad = 0.2
    add_text(ctx.slide, x + pad, y + 0.09, w - 2 * pad - 1.9,
             [{'text': b['label'], 'style': 'semi', 'size': 9, 'color': T.WHITE_70},
              {'text': b['name'], 'style': 'bold', 'size': 11.5, 'color': T.WHITE, 'before': 2}],
             rtl=ctx.rtl)
    bw = text_w(b['badge'], 'semi', 9) + 0.34
    bx = (x + pad) if ctx.rtl else (x + w - pad - bw)
    badge = rounded(ctx.slide, bx, y + 0.17, bw, 0.28, 0.14)
    solid(badge, T.WHITE); no_line(badge)
    tf = badge.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    if ctx.rtl:
        set_rtl_para(p)
    r = p.add_run(); r.text = b['badge']
    set_run_font(r, 'semi', 9, lc)
    return y + 0.62 + 0.16

def blk_ladder(ctx, b, x, y, w):
    """Stepped list (measurement ladder) with auto-height bars."""
    items = b['items']
    n = len(items)
    indent_total = b.get('indent', 1.6)
    base = _tick_color(b.get('color', 'teal'))
    pad = 0.14
    for i, it in enumerate(items):
        ind = indent_total * i / max(1, n - 1)
        bw = w - ind
        sx = x if ctx.rtl else (x + ind)
        tw = bw - 2 * pad - 0.30
        th = para_h(it, 'med', 9, tw, 1.30)
        bar_h = max(0.30, th + 0.13)
        sp = rounded(ctx.slide, sx, y, bw, bar_h, 0.07)
        tint = 0.16 + 0.14 * i
        mix = tuple(int(255 - (255 - c) * tint) for c in base)
        solid(sp, mix); no_line(sp)
        nx = (sx + bw - pad - 0.22) if ctx.rtl else (sx + pad)
        add_text(ctx.slide, nx, y + 0.055, 0.25,
                 [{'text': str(i + 1), 'style': 'bold', 'size': 11, 'color': T.NAVY, 'align': 'c'}])
        tx = (sx + pad) if ctx.rtl else (sx + pad + 0.30)
        add_text(ctx.slide, tx, y + 0.065, tw,
                 [{'text': it, 'style': 'med', 'size': 9, 'color': T.INK, 'spacing': 1.30}], rtl=ctx.rtl)
        y += bar_h + 0.055
    return y + 0.08

def blk_toc(ctx, b, x, y, w):
    for it in b['items']:
        num_w = 0.62
        nx = (x + w - num_w) if ctx.rtl else x
        add_text(ctx.slide, nx, y, num_w,
                 [{'text': it['num'], 'style': 'bold', 'size': 16, 'color': T.PURPLE,
                   'align': 'r' if ctx.rtl else 'l'}])
        tx = x + 0.5 if ctx.rtl else x + num_w + 0.12
        tw = w - num_w - 0.72
        th = add_text(ctx.slide, tx, y, tw,
                      [{'text': it['title'], 'style': 'semi', 'size': 11.5, 'color': T.INK}],
                      rtl=ctx.rtl)
        dh = add_text(ctx.slide, tx, y + th + 0.03, tw,
                      [{'text': it['desc'], 'style': 'reg', 'size': 8.5, 'color': T.GRAY,
                        'spacing': 1.35}], rtl=ctx.rtl)
        px = x if ctx.rtl else x + w - 0.4
        add_text(ctx.slide, px, y + 0.02, 0.4,
                 [{'text': str(it['page']), 'style': 'semi', 'size': 10, 'color': T.TEAL,
                   'align': 'l' if ctx.rtl else 'r'}])
        y += th + dh + 0.12
        ln = rect(ctx.slide, x, y, w, 0.008)
        solid(ln, T.CARD_BORDER); no_line(ln)
        y += 0.14
    return y

def blk_refs(ctx, b, x, y, w):
    tag_colors = b['tag_colors']
    for e in b['entries']:
        chip_w = 0.34
        nx = (x + w - chip_w) if ctx.rtl else x
        chip = rounded(ctx.slide, nx, y + 0.01, chip_w, 0.19, 0.05)
        solid(chip, tag_colors.get(e['tag'], T.GRAY)); no_line(chip)
        tf = chip.text_frame; tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = e['tag']
        set_run_font(r, 'semi', 8, T.WHITE)
        tx = x if ctx.rtl else x + chip_w + 0.10
        tw = w - chip_w - 0.10
        paras = [{'runs': [
            {'text': '[%d] ' % e['n'], 'style': 'semi', 'size': 8, 'color': T.PURPLE},
            {'text': e['text'], 'style': 'reg', 'size': 8, 'color': T.BODY}],
            'spacing': 1.30}]
        joined = '[%d] %s' % (e['n'], e['text'])
        th = para_h(joined, 'reg', 8, tw, 1.30)
        add_text(ctx.slide, tx, y, tw, paras, rtl=ctx.rtl)
        y += th + 0.015
        if e.get('url'):
            uh = add_text(ctx.slide, tx, y, tw,
                          [{'text': e['url'], 'style': 'reg', 'size': 7, 'color': T.TEAL, 'spacing': 1.22,
                            'link': e['url'], 'align': 'r' if ctx.rtl else 'l', 'ltr': True}],
                          rtl=False)
            y += uh + 0.02
        y += b.get('gap', 0.075)
    return y

BLOCKS = {
    'kicker': blk_kicker, 'h1': blk_h1, 'lede': blk_lede, 'body': blk_body,
    'note': blk_note, 'gap': blk_gap, 'subhead': blk_subhead, 'cards': blk_cards,
    'stats': blk_stats, 'panel': blk_panel, 'gpanel': blk_gpanel,
    'transition': blk_transition, 'findings': blk_findings, 'bullets': blk_bullets,
    'chips': blk_chips, 'pill': blk_pill, 'table': blk_table, 'image': blk_image,
    'columns': blk_columns, 'levelband': blk_levelband, 'ladder': blk_ladder,
    'toc': blk_toc, 'refs': blk_refs,
}

def render_blocks(ctx, blocks, x, y, w):
    for b in blocks:
        y = BLOCKS[b['t']](ctx, b, x, y, w)
    return y

# ---------------------------------------------------------------- pages
def page_cover(prs, meta, rtl, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(assets['cover_bg'], 0, 0, Inches(T.PAGE_W), Inches(T.PAGE_H))
    ctx = Ctx(slide, rtl, [], 1)
    # badge pill
    badge = meta['badge']
    bw = text_w(badge, 'semi', 9.5, 1.4) + 0.44
    bx = (T.PAGE_W - T.MARGIN - bw) if rtl else T.MARGIN
    pill = rounded(slide, bx, 4.35, bw, 0.34, 0.06)
    solid(pill, (0x1A, 0x14, 0x45)); no_line(pill)
    tf = pill.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    if rtl:
        set_rtl_para(p)
    r = p.add_run(); r.text = badge
    set_run_font(r, 'semi', 9.5, T.GREEN)
    if not rtl:
        set_spc(r, 1.4)
    # title
    y = 4.95
    for line, colr in meta['title_lines']:
        h = add_text(slide, T.MARGIN, y, T.CONTENT_W,
                     [{'text': line, 'style': 'bold', 'size': 30, 'color': colr, 'spacing': 1.18}],
                     rtl=rtl)
        y += h + 0.02
    y += 0.18
    add_text(slide, T.MARGIN, y, T.CONTENT_W,
             [{'text': meta['subtitle'], 'style': 'reg', 'size': 11.5, 'color': (0xE8, 0xE6, 0xF6),
               'spacing': 1.45}], rtl=rtl)
    # meta strip
    my = 9.55
    ln = rect(slide, T.MARGIN, my, T.CONTENT_W, 0.01)
    solid(ln, (0x4A, 0x42, 0x8A)); no_line(ln)
    cols = meta['meta_cols']
    cw = T.CONTENT_W / len(cols)
    for i, (label, val) in enumerate(cols):
        ci = (len(cols) - 1 - i) if rtl else i
        cx = T.MARGIN + ci * cw
        add_text(slide, cx, my + 0.16, cw - 0.2,
                 [{'text': label, 'style': 'semi', 'size': 8, 'color': (0x9E, 0x97, 0xC8),
                   'spc': 0 if rtl else 1.2},
                  {'text': val, 'style': 'reg', 'size': 10, 'color': T.WHITE, 'before': 4}],
                 rtl=rtl)
    return []

def page_section(prs, spec, rtl, page_no, doc_title, section_count=6):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ctx = Ctx(slide, rtl, [], page_no)
    x = T.MARGIN
    spc_lbl = 0 if rtl else 1.6
    title_h = para_h(spec['title'], 'bold', 22, T.CONTENT_W, 1.18)
    lede_h = para_h(spec['lede'], 'reg', 9.5, T.CONTENT_W - 1.0, 1.42)
    band_h = 0.42 + 0.24 + 0.10 + title_h + 0.10 + lede_h + 0.42
    band = rect(slide, 0, 0, T.PAGE_W, band_h)
    gradient(band, T.PURPLE, T.NAVY, 15)
    no_line(band)
    add_text(slide, x, 0.42, T.CONTENT_W,
             [{'text': spec['sec_label'], 'style': 'semi', 'size': 11, 'color': (0xCE, 0xC6, 0xEC),
               'spc': spc_lbl}], rtl=rtl)
    add_text(slide, x, 0.42 + 0.24 + 0.10, T.CONTENT_W,
             [{'text': spec['title'], 'style': 'bold', 'size': 22, 'color': T.WHITE,
               'spacing': 1.18}], rtl=rtl)
    ly = 0.42 + 0.24 + 0.10 + title_h + 0.10
    add_text(slide, x, ly, T.CONTENT_W - 1.0,
             [{'text': spec['lede'], 'style': 'reg', 'size': 9.5, 'color': (0xDD, 0xDA, 0xF0),
               'spacing': 1.42}], rtl=rtl)
    dw, dgap = 0.30, 0.10
    total_w = section_count * dw + (section_count - 1) * dgap
    dx = (T.PAGE_W - T.MARGIN - total_w) if rtl else x
    dy = band_h - 0.24
    for i in range(section_count):
        sp = rect(slide, dx + i * (dw + dgap), dy, dw, 0.05)
        solid(sp, T.GREEN if (i + 1) == spec['num'] else
              ((0xBF, 0xB8, 0xE2) if (i + 1) < spec['num'] else (0x6A, 0x5F, 0xA8)))
        no_line(sp)
    y = band_h + 0.30
    y = render_blocks(ctx, spec.get('blocks', []), x, y, T.CONTENT_W)
    footer(slide, page_no, doc_title, rtl)
    if y > T.FOOTER_RULE_Y - 0.05:
        ctx.warns.append((page_no, y))
    return ctx.warns

def page_content(prs, spec, rtl, page_no, doc_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ctx = Ctx(slide, rtl, [], page_no)
    y = render_blocks(ctx, spec['blocks'], T.MARGIN, T.TOP, T.CONTENT_W)
    footer(slide, page_no, doc_title, rtl)
    if y > T.FOOTER_RULE_Y - 0.05:
        ctx.warns.append((page_no, y))
    return ctx.warns

# ---------------------------------------------------------------- build deck
def build(pages, rtl, doc_title, assets, out_path):
    prs = Presentation()
    prs.slide_width = Inches(T.PAGE_W)
    prs.slide_height = Inches(T.PAGE_H)
    warns = []
    page_no = 0
    for spec in pages:
        page_no += 1
        if spec['type'] == 'cover':
            page_cover(prs, spec, rtl, assets)
        elif spec['type'] == 'section':
            warns += page_section(prs, spec, rtl, page_no, doc_title)
        else:
            warns += page_content(prs, spec, rtl, page_no, doc_title)
    prs.save(out_path)
    return warns
